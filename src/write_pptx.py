from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


def write_pptx(content: str, path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    slides_content = content.split("---")
    for slide_text in slides_content:
        lines = [line for line in slide_text.strip().split("\n") if line.strip()]
        if not lines:
            continue
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = lines[0]
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        for i, line in enumerate(lines[1:]):
            if i == 0:
                tf.text = line
            else:
                tf.add_paragraph().text = line
    prs.save(str(path))
